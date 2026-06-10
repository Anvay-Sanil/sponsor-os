"""Generate templates/master_deck.pptx — run once, commit the output.

The named-shape contract is the deliverable: core/deck_render.py fills shapes
by NAME, so a design-minded member can restyle this file in PowerPoint later
(fonts, positions, decorations) and the code keeps working, as long as the
shape names below survive.

Contract (shape names per slide):
  1 title    : accent_band, title_line, subtitle, test_banner
  2 story    : accent_bar, heading, hook, story
  3 numbers  : accent_bar, heading, fact_1..fact_4
  4 why_us   : accent_bar, heading, bullets, proof_note
  5 tiers    : accent_bar, heading, table_anchor, pricing_note
  6 cta      : accent_band, cta_line, contact
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "templates" / "master_deck.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
BLANK_LAYOUT = 6  # built-in blank layout index


def _box(slide, name: str, left: float, top: float, width: float, height: float,
         placeholder: str = "") -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    frame = shape.text_frame
    frame.word_wrap = True
    frame.text = placeholder
    frame.paragraphs[0].font.size = Pt(14)


def _band(slide, name: str, left: float, top: float, width: float, height: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.name = name
    shape.line.fill.background()


def build_template(path: Path = OUTPUT) -> Path:
    """Create the master deck template with the named-shape contract."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    layout = prs.slide_layouts[BLANK_LAYOUT]

    s1 = prs.slides.add_slide(layout)  # 1 — title
    _band(s1, "accent_band", 0, 0, 13.333, 2.6)
    _box(s1, "title_line", 0.8, 0.7, 11.7, 1.2, "{title_line}")
    _box(s1, "subtitle", 0.8, 3.1, 11.7, 0.9, "{subtitle}")
    _box(s1, "test_banner", 0.8, 6.3, 11.7, 0.8)

    s2 = prs.slides.add_slide(layout)  # 2 — story
    _band(s2, "accent_bar", 0, 0, 0.25, 7.5)
    _box(s2, "heading", 0.8, 0.5, 11.7, 0.9, "{heading}")
    _box(s2, "hook", 0.8, 1.7, 11.7, 1.2, "{hook}")
    _box(s2, "story", 0.8, 3.1, 11.7, 3.6, "{story}")

    s3 = prs.slides.add_slide(layout)  # 3 — fest in numbers
    _band(s3, "accent_bar", 0, 0, 0.25, 7.5)
    _box(s3, "heading", 0.8, 0.5, 11.7, 0.9, "{heading}")
    _box(s3, "fact_1", 0.8, 2.0, 5.6, 2.0, "{fact_1}")
    _box(s3, "fact_2", 6.9, 2.0, 5.6, 2.0, "{fact_2}")
    _box(s3, "fact_3", 0.8, 4.4, 5.6, 2.0, "{fact_3}")
    _box(s3, "fact_4", 6.9, 4.4, 5.6, 2.0, "{fact_4}")

    s4 = prs.slides.add_slide(layout)  # 4 — why us (evidence-anchored)
    _band(s4, "accent_bar", 0, 0, 0.25, 7.5)
    _box(s4, "heading", 0.8, 0.5, 11.7, 0.9, "{heading}")
    _box(s4, "bullets", 0.8, 1.8, 11.7, 4.4, "{bullets}")
    _box(s4, "proof_note", 0.8, 6.5, 11.7, 0.6, "{proof_note}")

    s5 = prs.slides.add_slide(layout)  # 5 — tiers & pricing
    _band(s5, "accent_bar", 0, 0, 0.25, 7.5)
    _box(s5, "heading", 0.8, 0.5, 11.7, 0.9, "{heading}")
    _box(s5, "table_anchor", 0.8, 1.8, 11.7, 4.2)
    _box(s5, "pricing_note", 0.8, 6.4, 11.7, 0.7, "{pricing_note}")

    s6 = prs.slides.add_slide(layout)  # 6 — call to action
    _band(s6, "accent_band", 0, 2.4, 13.333, 2.2)
    _box(s6, "cta_line", 0.8, 2.7, 11.7, 1.6, "{cta_line}")
    _box(s6, "contact", 0.8, 5.2, 11.7, 1.2, "{contact}")

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


if __name__ == "__main__":
    print(f"Wrote {build_template()}")
