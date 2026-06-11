"""Tests for Phase 3: facts gate, contrast guard, evidence anchoring, render."""
from __future__ import annotations

import io

import pytest
from PIL import Image
from pptx import Presentation

import core.chapter_facts as chapter_facts
from core.deck_render import TEMPLATE_PATH, TEST_BANNER_TEXT, render_deck
from core.palette import (
    DARK_TEXT,
    DEFAULT_PRIMARY,
    LIGHT_TEXT,
    contrast_ratio,
    dominant_color,
    ensure_contrast,
    readable_text_color,
)
from core.pitch import (
    DeckNarrative,
    EvidenceBullet,
    assemble_email,
    deck_path,
    fallback_bullets,
    next_deck_version,
    suggest_tier,
    valid_bullets,
)

# --- Catch 1: placeholder hard-block ------------------------------------------
def test_default_facts_have_placeholders_and_are_detected() -> None:
    assert chapter_facts.missing_facts(), "shipping file must list its own gaps"


def test_filled_facts_pass(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(chapter_facts, "CHAPTER_FACTS",
                        {key: "real value" for key in chapter_facts.CHAPTER_FACTS})
    assert chapter_facts.missing_facts() == []


# --- Catch 2: contrast guard ----------------------------------------------------
def test_black_on_white_is_max_contrast() -> None:
    assert contrast_ratio("#FFFFFF", "#000000") == pytest.approx(21.0)


def test_readable_text_on_light_and_dark() -> None:
    assert readable_text_color("#FFFF00") == DARK_TEXT   # yellow band -> dark text
    assert readable_text_color("#1A1A2E") == LIGHT_TEXT  # dark band -> white text


def test_unreadable_brand_color_falls_back_per_element() -> None:
    assert ensure_contrast("#FFFF66", "#FFFFFF") == DEFAULT_PRIMARY  # pale yellow on white
    assert ensure_contrast("#1A1A2E", "#FFFFFF") == "#1A1A2E"        # dark stays


def test_invalid_hex_falls_back() -> None:
    assert ensure_contrast("not-a-color", "#FFFFFF") == DEFAULT_PRIMARY


def test_dominant_color_of_red_image() -> None:
    buffer = io.BytesIO()
    Image.new("RGB", (32, 32), (220, 30, 40)).save(buffer, format="PNG")
    color = dominant_color(buffer.getvalue())
    assert color is not None
    red, green = int(color[1:3], 16), int(color[3:5], 16)
    assert red > 180 and green < 80


def test_dominant_color_ignores_white_images() -> None:
    buffer = io.BytesIO()
    Image.new("RGB", (32, 32), (255, 255, 255)).save(buffer, format="PNG")
    assert dominant_color(buffer.getvalue()) is None


def test_verbose_research_is_truncated_not_rejected() -> None:
    from core.pitch import BrandResearch

    research = BrandResearch(summary="x" * 5000, audience="y" * 5000,
                             india_activity="z" * 5000, tone="t" * 5000)
    assert len(research.summary) <= 600  # no quota-burning validation retry


def test_overlong_narrative_fields_clamp_instead_of_reject() -> None:
    narrative = DeckNarrative(
        title_line="boAt × ACM SIGAI MUJ: " + "x" * 200,  # observed live: placeholder
        hook="h" * 999,                                    # fest name blew the 60 cap
        story="s" * 5000,
        why_us_bullets=[EvidenceBullet(text="b" * 999, evidence_id=1)],
        cta_line="c" * 999,
        email_subject="e" * 999,
        email_body="m" * 5000,
    )
    assert len(narrative.title_line) == 60
    assert len(narrative.why_us_bullets[0].text) == 160
    assert len(narrative.email_body) == 1500


def test_too_short_output_still_rejects() -> None:
    with pytest.raises(Exception):
        DeckNarrative(title_line="x", hook="short", story="tiny",
                      why_us_bullets=[], cta_line="", email_subject="", email_body="")


# --- Catch 3: evidence anchoring --------------------------------------------------
def test_bullets_citing_unknown_evidence_are_dropped() -> None:
    bullets = [EvidenceBullet(text="Sponsored a rival fest recently", evidence_id=1),
               EvidenceBullet(text="Completely invented assertion here", evidence_id=999)]
    assert [b.evidence_id for b in valid_bullets(bullets, {1, 2})] == [1]


def test_fallback_bullets_are_verbatim_and_anchored() -> None:
    evidence = [{"id": 7, "source_type": "news",
                 "snippet": "[DEMO DATA] boAt backed a student hackathon."}]
    bullets = fallback_bullets(evidence)
    assert bullets[0].evidence_id == 7
    assert "boAt backed a student hackathon" in bullets[0].text
    assert "[DEMO DATA]" not in bullets[0].text


# --- versioning, tiers, email ------------------------------------------------------
def test_deck_paths_are_versioned_never_overwritten() -> None:
    assert deck_path(7, 1) == "7/deck_v1.pptx"
    assert deck_path(7, 4) == "7/deck_v4.pptx"


def test_version_counter_clears_storage_orphans() -> None:
    assert next_deck_version(0, []) == 1
    assert next_deck_version(3, []) == 4
    # orphan in storage with no table row (observed live 2026-06-11):
    assert next_deck_version(0, ["deck_v1.pptx"]) == 2
    assert next_deck_version(2, ["deck_v5.pptx", "unrelated.txt"]) == 6


@pytest.mark.parametrize(("score", "tier"),
                         [(92, "Gold"), (70, "Gold"), (69.9, "Silver"), (45, "Silver"),
                          (44.9, "Community"), (0, "Community")])
def test_tier_suggestion_rule(score: float, tier: str) -> None:
    assert suggest_tier(score) == tier


def _narrative() -> DeckNarrative:
    return DeckNarrative(
        title_line="boAt × Demo Fest 2026",
        hook="Your audience is already in our crowd — here is the proof.",
        story="Three thousand students, two days, and a brand wall your reels team "
              "will love. This is a working story paragraph for the smoke test.",
        why_us_bullets=[EvidenceBullet(text="Documented on a rival fest's website: "
                                            "boAt sponsored their main stage.", evidence_id=1)],
        cta_line="Let's talk for 15 minutes this week.",
        email_subject="Partnership idea: boAt × Demo Fest",
        email_body="Hi team — we run the ACM SIGAI chapter fest at MUJ and noticed "
                   "boAt backing fests like ours. We have a package that fits your "
                   "campus push and real numbers to back it. Could we take 15 minutes "
                   "this week to walk you through it? Happy to send the deck ahead.",
    )


def test_email_test_prefix() -> None:
    assert assemble_email(_narrative(), is_test=True).startswith("[TEST — DO NOT SEND]")
    assert not assemble_email(_narrative(), is_test=False).startswith("[TEST")


# --- render smoke (no LLM, no network) ----------------------------------------------
def test_template_exists() -> None:
    assert TEMPLATE_PATH.exists(), "run: python jobs/build_master_template.py"


def _all_deck_text(pptx_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(pptx_bytes))
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return " ".join(chunks)


TIERS = [{"name": "Gold", "base_price": 50000,
          "components_json": {"stage_logo": True, "reels": 2}}]


def test_render_real_deck_smoke() -> None:
    pptx_bytes = render_deck(_narrative(), "boAt", {"primary": "#FF0040"}, TIERS,
                             "Gold", chapter_facts.CHAPTER_FACTS, is_test=False)
    assert pptx_bytes[:2] == b"PK"
    text = _all_deck_text(pptx_bytes)
    assert "boAt × Demo Fest 2026" in text
    assert "Indicative pricing" in text
    assert "₹50,000" in text
    assert TEST_BANNER_TEXT not in text


def test_render_test_deck_carries_watermark() -> None:
    pptx_bytes = render_deck(_narrative(), "boAt", {"primary": "#FF0040"}, TIERS,
                             "Gold", chapter_facts.CHAPTER_FACTS, is_test=True)
    assert TEST_BANNER_TEXT in _all_deck_text(pptx_bytes)
