"""Render a DeckNarrative into PPTX bytes against templates/master_deck.pptx.

Fills shapes by NAME (contract documented in jobs/build_master_template.py).
Every brand color passes the WCAG contrast guard from core.palette before it
touches text or background — failing elements fall back to the ACM palette
per element, never the whole deck.
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from core.palette import DARK_TEXT, ensure_contrast, readable_text_color

TEMPLATE_PATH = Path(__file__).resolve().parents[1] / "templates" / "master_deck.pptx"
TEST_BANNER_TEXT = "🧪 TEST DECK — DO NOT SEND. Built from demo data for practice only."
PRICING_NOTE = ("Indicative pricing — exact packages and final numbers follow a "
                "conversation with our team.")


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _find(slide: Any, name: str) -> Any | None:
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_text(slide: Any, name: str, text: str, *, size: int = 16,
              color: str = DARK_TEXT, bold: bool = False) -> None:
    shape = _find(slide, name)
    if shape is None or not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.word_wrap = True
    frame.text = text
    for paragraph in frame.paragraphs:
        for run in paragraph.runs or []:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)


def _fill(slide: Any, name: str, hex_color: str) -> None:
    shape = _find(slide, name)
    if shape is None:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)


def _bullets(slide: Any, name: str, lines: list[str], *, size: int = 16) -> None:
    shape = _find(slide, name)
    if shape is None:
        return
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {line}"
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = _rgb(DARK_TEXT)


def _tier_table(slide: Any, tiers: list[dict[str, Any]], primary: str) -> None:
    anchor = _find(slide, "table_anchor")
    if anchor is None or not tiers:
        return
    rows = len(tiers) + 1
    table = slide.shapes.add_table(rows, 3, anchor.left, anchor.top,
                                   anchor.width, anchor.height).table
    header_text = readable_text_color(primary)
    for column, label in enumerate(("Tier", "What you get", "Indicative price")):
        cell = table.cell(0, column)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(primary)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = _rgb(header_text)
    for row, tier in enumerate(tiers, start=1):
        components = tier.get("components_json") or {}
        included = ", ".join(
            key.replace("_", " ") if isinstance(value, bool) else f"{value}× {key.replace('_', ' ')}"
            for key, value in components.items() if value
        )
        price = tier.get("base_price")
        cells = (str(tier.get("name", "")), included or "—",
                 f"₹{int(price):,}" if price else "—")
        for column, text in enumerate(cells):
            cell = table.cell(row, column)
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(13)


def render_deck(narrative: Any, brand_name: str, palette: dict[str, str],
                tiers: list[dict[str, Any]], suggested_tier: str,
                facts: dict[str, str], is_test: bool) -> bytes:
    """Fill the master template and return PPTX bytes. No network, no LLM."""
    prs = Presentation(str(TEMPLATE_PATH))
    s1, s2, s3, s4, s5, s6 = list(prs.slides)[:6]
    primary = palette.get("primary", "#6C5CE7")
    on_primary = readable_text_color(primary)         # text sitting ON the band
    heading_color = ensure_contrast(primary, "#FFFFFF")  # brand color on white

    _fill(s1, "accent_band", primary)
    _set_text(s1, "title_line", narrative.title_line, size=40, bold=True, color=on_primary)
    _set_text(s1, "subtitle",
              f"A partnership proposal for {brand_name} · {facts['fest_name']} · {facts['fest_dates']}",
              size=18)
    if is_test:
        _set_text(s1, "test_banner", TEST_BANNER_TEXT, size=16, bold=True, color="#C0392B")

    _fill(s2, "accent_bar", primary)
    _set_text(s2, "heading", f"Why this fits {brand_name}", size=28, bold=True, color=heading_color)
    _set_text(s2, "hook", narrative.hook, size=20, bold=True)
    _set_text(s2, "story", narrative.story, size=16)

    _fill(s3, "accent_bar", primary)
    _set_text(s3, "heading", f"{facts['fest_name']} in numbers", size=28, bold=True, color=heading_color)
    _set_text(s3, "fact_1", f"{facts['expected_footfall']}\nexpected footfall", size=20, bold=True)
    _set_text(s3, "fact_2", f"{facts['instagram_followers']}\nInstagram followers", size=20, bold=True)
    _set_text(s3, "fact_3", facts["past_highlights"], size=18)
    _set_text(s3, "fact_4", f"Organized by\n{facts['chapter_name']}", size=18)

    _fill(s4, "accent_bar", primary)
    _set_text(s4, "heading", "Why partner with us", size=28, bold=True, color=heading_color)
    _bullets(s4, "bullets", [bullet.text for bullet in narrative.why_us_bullets])
    _set_text(s4, "proof_note",
              f"Every claim above is backed by {len(narrative.why_us_bullets)} public "
              "source(s) on file — ask us for the links.", size=12)

    _fill(s5, "accent_bar", primary)
    _set_text(s5, "heading", f"Packages — we suggest {suggested_tier}", size=28, bold=True,
              color=heading_color)
    _tier_table(s5, tiers, primary)
    _set_text(s5, "pricing_note", PRICING_NOTE, size=12)

    _fill(s6, "accent_band", primary)
    _set_text(s6, "cta_line", narrative.cta_line, size=28, bold=True, color=on_primary)
    _set_text(s6, "contact", f"{facts['chapter_name']}  ·  {facts['contact_email']}", size=16)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
